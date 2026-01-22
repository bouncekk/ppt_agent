from __future__ import annotations

from dataclasses import dataclass, asdict
from pathlib import Path
from typing import List
import json

from pptx import Presentation


@dataclass
class Slide:
    index: int
    title: str
    bullets: List[str]
    notes: str | None = None


def parse_ppt(path: str | Path) -> List[Slide]:
    """使用 python-pptx 将 PPT 文件解析为 Slide 列表的最小可用实现。

    解析策略：
    - 遍历每一页 slide；
    - 收集当前页中所有带文本的 shape；
    - 第一段非空文本作为标题 `title`；
    - 其余非空文本行按段落归入 `bullets`；
    - 暂不区分备注区，`notes` 先置为 None，后续可扩展从 notes_slide 抽取。

    注意：本函数假定输入文件为 .pptx 格式。
    """

    ppt_path = Path(path)
    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT 文件不存在: {ppt_path}")

    presentation = Presentation(ppt_path)
    slides: List[Slide] = []

    for idx, slide in enumerate(presentation.slides, start=1):
        texts: List[str] = []

        # 收集当前页所有文本框中的段落文本
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            if shape.text_frame is None:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = (paragraph.text or "").strip()
                if text:
                    texts.append(text)

        if texts:
            title = texts[0]
            bullets = texts[1:] if len(texts) > 1 else []
        else:
            # 空白页的兜底处理
            title = f"Slide {idx}"
            bullets = []

        notes: str | None = None

        slides.append(
            Slide(
                index=idx,
                title=title,
                bullets=bullets,
                notes=notes,
            )
        )

    return slides


def slides_to_json(slides: List[Slide]) -> str:
    """将 Slide 列表序列化为 JSON 字符串。"""

    return json.dumps([asdict(s) for s in slides], ensure_ascii=False, indent=2)


def parse_ppt_to_json_file(ppt_path: str | Path, output_path: str | Path) -> Path:
    """占位管线：PPT → Slide 结构 → JSON 文件。
    """

    slides = parse_ppt(ppt_path)
    output_path = Path(output_path)
    output_path.write_text(slides_to_json(slides), encoding="utf-8")
    return output_path
